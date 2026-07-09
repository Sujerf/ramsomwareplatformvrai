#!/usr/bin/env python3
"""Met à jour les notes orateur du PPTX avec le texte allégé."""
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

NOTES = {
    1: """[OUVERTURE]

Excellence Monsieur le Président du jury, Mesdames et Messieurs les membres du jury, bonjour.

Permettez-nous d'exprimer notre profonde reconnaissance pour votre présence. Ce jour est l'aboutissement de ces années de formation, un moment qui restera gravé dans notre parcours.

Le sujet de notre étude est intitulé : Conception et mise en place d'un système de détection et de réponse aux attaques ransomware dans un environnement contrôlé.

C'est avec fierté et humilité que nous vous présentons RansomShield — un système conçu de bout en bout pour détecter et contrer les attaques ransomware, de manière contrôlée et traçable, avec une validation humaine à chaque étape décisive.

→ Voici comment s'articule cette présentation.""",

    2: """[SOMMAIRE]

Notre présentation s'articule en cinq parties. Nous partirons du contexte et de la problématique, pour aller vers un état des lieux des solutions existantes, puis la conception et la réalisation. Nous conclurons par les perspectives et terminerons par une démonstration en conditions réelles.

→ Commençons par le contexte et la problématique.""",

    4: """[CONTEXTE]

Aujourd'hui, les entreprises, les hôpitaux, les universités fonctionnent grâce à leurs données numériques. Le ransomware — un logiciel malveillant qui rend les données inaccessibles puis exige une rançon — exploite précisément cette dépendance. En 2023, ces attaques ont progressé de 73 % selon l'ENISA, avec un coût moyen de plus d'un million de dollars par incident. Ce qui est grave : sans outil adapté, une attaque passe en moyenne 21 jours sans être détectée.

→ Ce constat nous amène à la problématique centrale.""",

    5: """[PROBLÉMATIQUE]

Les antivirus fonctionnent par reconnaissance d'empreintes : ils connaissent les menaces répertoriées. Mais face à une variante inédite, ils restent aveugles — et pendant ce temps, les données sont chiffrées, les sauvegardes effacées, les dommages deviennent irréversibles.

Il n'existe pas de solution libre et démontrable pour détecter ces attaques à partir de leurs comportements, avec un contrôle humain intégré à chaque étape.

Trois défis ont structuré ce travail : les variantes inconnues, la diversité des systèmes à surveiller, et le risque qu'une réponse automatique non validée cause elle-même des dommages.

→ Ces défis ont défini les six objectifs que voici.""",

    6: """[OBJECTIFS]

Six objectifs précis ont guidé ce travail : surveiller en temps réel les fichiers sur chaque terminal ; calculer un niveau de risque ; déclencher des alertes automatiquement. Le principe fondateur : chaque action passe par la décision d'un opérateur humain. L'accès est sécurisé par double authentification — mot de passe et code temporaire valide 30 secondes. Et pour valider l'ensemble, des scénarios d'attaque représentatifs ont été définis et exécutés.

→ Voyons maintenant les fondements théoriques de ce travail.""",

    8: """[DÉFINITION]

Un ransomware agit toujours selon le même schéma : il renomme vos fichiers avec une extension inhabituelle pour les rendre illisibles, dépose une note de rançon, puis efface les sauvegardes pour rendre toute récupération impossible. Il utilise parfois des outils légitimes pour passer inaperçu.

Ce sont ces comportements observables — ces traces — qui forment la base de notre approche. Surveiller ces actions en temps réel permet de réagir avant que les dommages soient irréversibles.

→ Mais pourquoi les outils existants ne suffisent-ils pas ?""",

    9: """[SOLUTIONS EXISTANTES]

Des solutions existent — voyons pourquoi elles restent insuffisantes.

Les antivirus ne reconnaissent que les menaces déjà cataloguées. Les SIEM collectent des informations, mais sans chercher activement des comportements suspects. Les EDR sont très efficaces, mais leur coût les rend inaccessibles pour la majorité des organisations — PME, hôpitaux, administrations — qui n'ont pas les budgets des grands groupes.

RansomShield occupe cet espace laissé vacant : une solution libre, comportementale, déployable sur n'importe quelle infrastructure, avec un contrôle humain à chaque étape.

→ Voyons maintenant comment ce système a été conçu.""",

    11: """[ARCHITECTURE]

Voici l'architecture générale — à lire de gauche à droite.

Un agent sur chaque terminal observe toute activité inhabituelle sur les fichiers. Ces observations sont transmises au serveur central, qui calcule un score de risque. Dès que ce score dépasse le seuil configuré, une alerte est générée et une action de protection est proposée à l'opérateur. C'est lui qui prend la décision finale. Le système n'agit jamais seul.

→ Voyons de plus près les quatre composants.""",

    12: """[4 COMPOSANTS]

L'agent de surveillance, d'abord : installé sur chaque machine, il observe en temps réel toute activité inhabituelle sur les fichiers et conserve ces informations localement, pour ne rien perdre même en cas de coupure réseau.

La console de supervision, ensuite : c'est le poste de commandement, là où l'opérateur suit les alertes, prend ses décisions et consulte l'historique complet de toutes les actions.

Le module de découverte réseau : il signale automatiquement les nouvelles machines, mais aucune n'est intégrée au système sans accord explicite de l'administrateur.

Enfin, l'outil de validation : il rejoue les comportements documentés des ransomwares — renommage de fichiers, dépôt de note de rançon, suppression des sauvegardes — sans déployer de malware réel. C'est une contrainte éthique et légale du cadre académique ; les comportements simulés restent fidèles aux cas réels documentés par l'ENISA.

→ Qui sont les acteurs qui interagissent avec ce système ?""",

    13: """[USE CASE]

Deux profils d'utilisateurs interagissent avec le système. L'administrateur configure et pilote : il connecte les agents, définit les règles et lance les simulations. L'analyste SOC traite les incidents au quotidien : il consulte les alertes, approuve ou rejette les actions proposées, et suit l'historique des événements. L'analyste ne peut pas modifier la configuration — cette séparation garantit l'intégrité du système.

→ Passons maintenant à la réalisation concrète.""",

    15: """[CHOIX TECHNIQUES]

Les choix techniques ont été guidés par trois critères : portabilité multi-OS, facilité de déploiement, et démontrable sans infrastructure coûteuse.

Le serveur et la console sont développés avec Laravel. L'agent de surveillance est en Python, portable sur Windows, Linux et macOS. Les tests ont été conduits sur trois machines virtuelles en réseau isolé. L'accès à la console est protégé par double authentification : mot de passe puis code temporaire valide 30 secondes.

→ Voyons concrètement comment l'agent et le moteur fonctionnent.""",

    16: """[AGENT & MOTEUR DE DÉTECTION]

L'agent démarre, est reconnu par le serveur grâce à un jeton sécurisé, surveille les dossiers et envoie un signal de vie toutes les 30 secondes.

Pour le moteur de scoring : un fichier renommé .locked ajoute 80 points au score de risque. Une note de rançon dans le même dossier en ajoute 55. Le score atteint 135 — au-dessus du seuil critique de 100. En moins de cinq secondes, une alerte est générée et un incident est ouvert dans la console, en attente de la décision de l'opérateur.

→ Côté analyste, voici ce que la console SOC offre.""",

    17: """[CONSOLE SOC]

Venons-en maintenant à la console SOC — l'interface que l'opérateur utilise au quotidien pour surveiller, réagir et rendre des comptes.

La console regroupe vingt-trois pages en quatre sections : supervision, infrastructure, détection et réponse, et configuration.

Ce qui est fondamental : chaque action proposée passe obligatoirement par la file d'approbation. L'analyste l'approuve ou la rejette explicitement. Aucune action ne s'exécute en silence.

Dès qu'une alerte est levée, l'opérateur est notifié par quatre canaux simultanément : interface web, e-mail, signal sonore et webhook.

→ Qu'est-ce que les tests ont donné ?""",

    18: """[TESTS ET VALIDATION]

Cinq scénarios d'attaque ont été définis, du chiffrement simple à la chaîne complète en 22 événements, exécutés sur trois machines virtuelles en réseau isolé — sans jamais déployer de malware réel. Pour chaque scénario, trois indicateurs ont été mesurés : délai de détection, intégrité des données, et complétude de la chaîne jusqu'à la décision de l'opérateur.

Les résultats sont constants : cinq scénarios exécutés, cinq détectés en moins de cinq secondes, zéro événement perdu. Limite assumée : les tests ont été conduits sur Linux uniquement.

→ Place maintenant à la démonstration live.""",

    19: """[DÉMONSTRATION LIVE]

Suivre le guide de démonstration.

1. Console SOC ouverte : 3 agents en ligne
2. SSH sur rs-client-1 : renommage .locked + note de rançon
3. Revenir sur la console : alerte critique en moins de 5 secondes
4. File d'approbation : approuver l'action, montrer le rollback
5. Timeline : tout est horodaté
6. Simulation kill chain complète : 22 événements

→ Revenir à la présentation pour conclure.""",

    21: """[CONCLUSION]

RansomShield couvre l'intégralité du cycle de traitement d'un incident : détection sur le terminal, analyse comportementale, alerte, décision de l'opérateur, archivage complet. Il est opérationnel, démontrable, construit sans malware réel.

Nous n'avons pas la prétention d'avoir épuisé le sujet : chaque remarque du jury est pour nous une opportunité de progresser et de mieux faire.

Les perspectives sont tracées : sécuriser les communications pour la production, affiner les rôles, et doter le moteur d'un apprentissage automatique pour reconnaître les anomalies avec plus de précision.

→ Nous restons à votre disposition pour vos questions.""",

    22: """[MERCI — MOT DE CLÔTURE]

Dans un monde où les systèmes automatisent nos décisions, RansomShield rappelle une conviction simple : la machine peut détecter, mais c'est l'humain qui décide. Ce n'est pas une contrainte technique — c'est un choix éthique.

Merci pour votre aimable attention et pour le temps accordé à ce travail.""",
}


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    # Supprimer tous les paragraphes sauf le premier
    for para in list(tf.paragraphs)[1:]:
        p = para._p
        p.getparent().remove(p)
    # Réécrire le premier paragraphe
    first_para = tf.paragraphs[0]
    p_elem = first_para._p
    # Supprimer tous les runs existants
    for r in p_elem.findall(qn('a:r')):
        p_elem.remove(r)
    # Créer un nouveau run avec le texte
    r_elem = etree.SubElement(p_elem, qn('a:r'))
    rPr = etree.SubElement(r_elem, qn('a:rPr'))
    rPr.set('lang', 'fr-FR')
    rPr.set('dirty', '0')
    t_elem = etree.SubElement(r_elem, qn('a:t'))
    t_elem.text = text


prs = Presentation('presentation_soutenance.pptx')

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    if slide_num in NOTES:
        set_notes(slide, NOTES[slide_num])
        print(f"Slide {slide_num} mise à jour.")

prs.save('presentation_soutenance.pptx')
print("PPTX sauvegardé.")
