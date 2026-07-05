#!/usr/bin/env python3
"""
Reconstruit le PPTX en remplaçant chaque slide par l'image PNG correspondante
tout en conservant les notes orateur existantes.
"""
import sys, os
sys.path.insert(0, '.venv-pptx/lib/python3.12/site-packages')

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy

PPTX    = "presentation_soutenance.pptx"
SLIDES  = ".pptx_slides"

prs = Presentation(PPTX)

slide_w = prs.slide_width
slide_h = prs.slide_height

for i, slide in enumerate(prs.slides, start=1):
    img_path = os.path.join(SLIDES, f"slide-{i:02d}.png")
    if not os.path.exists(img_path):
        print(f"  Slide {i} : image introuvable, ignorée")
        continue

    # Sauvegarder les notes avant de modifier la slide
    notes_text = ""
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text

    # Supprimer toutes les shapes existantes
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp = shape._element
        sp_tree.remove(sp)

    # Ajouter l'image pleine page
    slide.shapes.add_picture(img_path, 0, 0, slide_w, slide_h)

    # Restaurer les notes
    if notes_text.strip():
        tf = slide.notes_slide.notes_text_frame
        for para in tf.paragraphs[1:]:
            para._p.getparent().remove(para._p)
        tf.paragraphs[0].text = notes_text

    print(f"✓ Slide {i} mise à jour")

prs.save(PPTX)
print(f"\nOK : {PPTX} reconstruit avec les nouvelles images.")
