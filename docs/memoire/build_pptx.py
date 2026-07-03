#!/usr/bin/env python3
"""
Génère presentation_soutenance.pptx à partir du PDF propre (sans notes) et
des blocs \\note{...} du fichier .tex.

Chaque slide devient une image pleine page + les notes orateur natives
PowerPoint (View > Notes), invisibles automatiquement en mode Diaporama /
partage d'écran.

Prérequis : ./build.sh doit avoir été lancé avant (pour produire
presentation_soutenance.pdf), et un environnement Python avec python-pptx
et Pillow (voir requirements-pptx.txt).

Usage : python3 build_pptx.py
"""
import glob
import os
import re
import subprocess
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
TEX = os.path.join(HERE, "presentation_soutenance.tex")
PDF = os.path.join(HERE, "presentation_soutenance.pdf")
SLIDES_DIR = os.path.join(HERE, ".pptx_slides")
OUT = os.path.join(HERE, "presentation_soutenance.pptx")


def render_slide_images():
    if os.path.exists(SLIDES_DIR):
        for f in glob.glob(os.path.join(SLIDES_DIR, "*.png")):
            os.remove(f)
    else:
        os.makedirs(SLIDES_DIR)
    subprocess.run(
        ["pdftoppm", "-png", "-r", "200", PDF, os.path.join(SLIDES_DIR, "slide")],
        check=True,
    )


def strip_braced_command(text, marker):
    """Remove a \\newcommand{...}[n]{ BODY } block given its marker prefix."""
    idx = text.find(marker)
    if idx == -1:
        return text
    brace_start = text.index("{", idx + len(marker))
    depth = 1
    i = brace_start + 1
    while depth > 0:
        if text[i] == "{":
            depth += 1
        elif text[i] == "}":
            depth -= 1
        i += 1
    return text[:idx] + text[i:]


def extract_braced(text, start_after_brace):
    depth = 1
    i = start_after_brace
    buf = []
    while depth > 0:
        c = text[i]
        if c == "{":
            depth += 1
        elif c == "}":
            depth -= 1
            if depth == 0:
                break
        buf.append(c)
        i += 1
    return "".join(buf), i + 1


def parse_notes_by_frame():
    with open(TEX, encoding="utf-8") as f:
        src = f.read()

    clean_lines = [ln for ln in src.split("\n") if not ln.lstrip().startswith("%")]
    src_nc = "\n".join(clean_lines)

    # The \sectionslide macro DEFINITION contains a literal \begin{frame}
    # (template code, not an invocation) -- strip it before counting frames.
    src_nc = strip_braced_command(src_nc, r"\newcommand{\sectionslide}[3]")

    frame_re = re.compile(r"\\begin\{frame\}|\\sectionslide\{")
    note_re = re.compile(r"\\note\{")

    events = []
    for m in frame_re.finditer(src_nc):
        events.append((m.start(), "frame", None))
    for m in note_re.finditer(src_nc):
        events.append((m.start(), "note_start", m.end()))
    events.sort(key=lambda e: e[0])

    frame_index = 0
    notes_by_frame = {}
    for _, kind, payload in events:
        if kind == "frame":
            frame_index += 1
        else:
            content, _ = extract_braced(src_nc, payload)
            notes_by_frame[frame_index] = content

    return frame_index, notes_by_frame


def latex_to_text(s):
    s = s.strip()
    s = re.sub(r"\\textbf\{(\[[^\]]*\])\}", r"\1\n", s)
    for _ in range(3):
        for cmd in ["textbf", "textit", "texttt", "small", "textsuperscript"]:
            s = re.sub(rf"\\{cmd}\{{([^{{}}]*)\}}", r"\1", s)
    s = re.sub(r"\\og\s*", "« ", s)
    s = re.sub(r"\\fg\{\}", " »", s)
    s = re.sub(r"\\fg\s*", " »", s)
    s = s.replace(r"\textbullet\ \ ", "• ").replace(r"\textbullet\ ", "• ").replace(r"\textbullet", "•")
    s = s.replace(r"$\to$", "→")
    s = s.replace(r"\%", "%")
    s = s.replace(r"~", " ")
    s = s.replace(r"\,", " ")
    s = s.replace(r"\enspace", " ")
    s = s.replace(r"\quad", "  ")
    s = s.replace(r"\&", "&")
    s = s.replace("---", "—")
    s = s.replace("--", "–")
    s = re.sub(r"\\\\\[\d+pt\]", "\n", s)
    s = s.replace(r"\\", "\n")
    s = re.sub(r"\\[a-zA-Z]+\{([^{}]*)\}", r"\1", s)
    s = re.sub(r"\\[a-zA-Z]+", "", s)

    out_lines = [re.sub(r"[ \t]+", " ", line).strip() for line in s.split("\n")]
    result, prev_blank = [], False
    for line in out_lines:
        blank = line == ""
        if blank and prev_blank:
            continue
        result.append(line)
        prev_blank = blank
    return "\n".join(result).strip()


def main():
    if not os.path.exists(PDF):
        sys.exit("presentation_soutenance.pdf introuvable — lancez d'abord ./build.sh")

    render_slide_images()

    frame_count, notes_by_frame = parse_notes_by_frame()
    notes_text = {k: latex_to_text(v) for k, v in notes_by_frame.items()}

    img_files = sorted(glob.glob(os.path.join(SLIDES_DIR, "slide-*.png")))
    if len(img_files) != frame_count:
        sys.exit(
            f"Incohérence : {len(img_files)} images rendues vs {frame_count} "
            f"frames détectées dans le .tex — vérifier le parsing."
        )

    from pptx import Presentation
    from pptx.util import Inches, Emu
    from PIL import Image

    with Image.open(img_files[0]) as im:
        px_w, px_h = im.size
    aspect = px_w / px_h

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Emu(int(prs.slide_width / aspect))
    blank_layout = prs.slide_layouts[6]

    for idx, img_path in enumerate(img_files, start=1):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        note_txt = notes_text.get(idx, "")
        if note_txt:
            slide.notes_slide.notes_text_frame.text = note_txt

    prs.save(OUT)
    print(f"OK : {OUT} ({len(img_files)} slides, {len(notes_text)} avec notes)")


if __name__ == "__main__":
    main()
