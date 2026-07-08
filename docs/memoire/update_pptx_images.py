#!/usr/bin/env python3
"""
Met à jour les images du PPTX avec les pages PDF recompilées.
Mapping PPTX slide → PDF page (les slides 14 et 15 sont orphelines, on les garde).
"""
import os, shutil
from pptx import Presentation
from pptx.util import Emu
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
PPTX = os.path.join(HERE, "presentation_soutenance.pptx")
IMG_PREFIX = "/tmp/pptx_slides"   # pdftoppm -r 150 -png → /tmp/pptx_slides-01.png

# PPTX slide (1-based) → PDF page (1-based) — mapping 1-to-1 (22 slides = 22 pages)
MAPPING = {i: i for i in range(1, 23)}


def replace_image(slide, png_path):
    """Remplace l'image du premier Picture shape d'une slide."""
    pic_shape = next(s for s in slide.shapes if s.shape_type == 13)
    pic = pic_shape._element
    NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    blipFill = pic.find(f'{{{NS_P}}}blipFill')
    blip     = blipFill.find(f'{{{NS_A}}}blip')
    rId      = blip.get(f'{{{NS_R}}}embed')
    part = slide.part.related_part(rId)
    with open(png_path, 'rb') as f:
        part._blob = f.read()
    # content_type est déjà image/png (ou jpeg) — pas besoin de le changer


prs = Presentation(PPTX)
slides = list(prs.slides)

updated = 0
for pptx_idx, pdf_page in MAPPING.items():
    if pdf_page is None:
        print(f"  Slide {pptx_idx:2d}: conservée (orpheline)")
        continue
    img = f"{IMG_PREFIX}-{pdf_page:02d}.png"
    if not os.path.exists(img):
        print(f"  Slide {pptx_idx:2d}: IMAGE MANQUANTE → {img}")
        continue
    replace_image(slides[pptx_idx - 1], img)
    print(f"  Slide {pptx_idx:2d} ← PDF page {pdf_page:2d}  ✓")
    updated += 1

prs.save(PPTX)
print(f"\nOK : {updated} slides mises à jour → {PPTX}")
